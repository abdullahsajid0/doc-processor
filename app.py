import os
import streamlit as st
from groq import Groq
import pdfplumber
from pptx import Presentation
import pandas as pd
from docx import Document
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.units import inch
import time
from typing import Optional, Dict, List
import logging

# Configure custom theme and styling
st.set_page_config(
    page_title="Smart Document Processor",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Enhanced CSS with animations and professional styling
st.markdown("""
    <style>
        /* Global styles */
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        * {
            font-family: 'Inter', sans-serif;
        }

        [data-testid="stAppViewContainer"] {
            background: linear-gradient(135deg, #f5f7fa 0%, #e4e8eb 100%);
            padding: 1.5rem;
        }

        /* Main container styling */
        .main {
            max-width: 1200px;
            margin: 0 auto;
        }

        /* Enhanced title container */
        .title-container {
            background: linear-gradient(135deg, #2193b0 0%, #6dd5ed 100%);
            padding: 2rem;
            border-radius: 20px;
            margin-bottom: 2rem;
            box-shadow: 0 10px 30px rgba(0, 0, 0, 0.1);
            text-align: center;
            animation: fadeIn 0.8s ease-out;
        }

        .title-container h1 {
            font-size: 2.5rem;
            font-weight: 700;
            margin-bottom: 0.5rem;
            background: linear-gradient(120deg, #ffffff 0%, #f0f0f0 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }

        .title-container p {
            color: rgba(255, 255, 255, 0.9);
            font-size: 1.1rem;
            font-weight: 500;
        }

        /* Enhanced card styling */
        .stCard {
            background: white;
            border-radius: 20px;
            padding: 1.5rem;
            box-shadow: 0 8px 20px rgba(0, 0, 0, 0.06);
            margin: 1.2rem 0;
            transition: transform 0.3s ease, box-shadow 0.3s ease;
            animation: slideUp 0.5s ease-out;
        }

        .stCard:hover {
            transform: translateY(-5px);
            box-shadow: 0 12px 30px rgba(0, 0, 0, 0.1);
        }

        /* File uploader enhancements */
        .uploadedFile {
            background: rgba(33, 147, 176, 0.05);
            border-radius: 15px;
            padding: 1rem;
            margin: 1rem 0;
            border: 2px dashed #2193b0;
            transition: all 0.3s ease;
        }

        .uploadedFile:hover {
            background: rgba(33, 147, 176, 0.1);
            border-color: #6dd5ed;
        }

        /* Enhanced metric cards */
        .metric-container {
            background: white;
            padding: 1.2rem;
            border-radius: 15px;
            text-align: center;
            box-shadow: 0 4px 15px rgba(0, 0, 0, 0.05);
            transition: transform 0.3s ease;
            animation: fadeIn 0.5s ease-out;
        }

        .metric-container:hover {
            transform: translateY(-3px);
        }

        .metric-value {
            font-size: 2rem;
            font-weight: 700;
            color: #2193b0;
            margin-bottom: 0.3rem;
        }

        .metric-label {
            color: #666;
            font-size: 0.9rem;
            font-weight: 500;
        }

        /* Enhanced buttons */
        .stButton>button {
            background: linear-gradient(135deg, #2193b0 0%, #6dd5ed 100%);
            color: white;
            border-radius: 12px;
            padding: 0.8rem 1.5rem;
            font-weight: 600;
            border: none;
            transition: all 0.3s ease;
            box-shadow: 0 4px 15px rgba(33, 147, 176, 0.2);
            width: auto;
            min-width: 150px;
        }

        .stButton>button:hover {
            transform: translateY(-2px);
            box-shadow: 0 6px 20px rgba(33, 147, 176, 0.3);
        }

        /* Download button special styling */
        .download-button {
            background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%) !important;
        }

        /* Enhanced response container */
        .response-container {
            background: white;
            padding: 1.5rem;
            border-radius: 15px;
            margin: 1.2rem 0;
            border-left: 5px solid #2193b0;
            box-shadow: 0 4px 15px rgba(0, 0, 0, 0.05);
            animation: slideIn 0.5s ease-out;
        }

        .response-container h4 {
            color: #2193b0;
            font-size: 1.2rem;
            margin-bottom: 1rem;
        }

        /* Radio button enhancements */
        .stRadio > div {
            background: white;
            padding: 1rem;
            border-radius: 12px;
            box-shadow: 0 4px 15px rgba(0, 0, 0, 0.05);
        }

        /* Progress bar enhancement */
        .stProgress > div > div {
            background: linear-gradient(135deg, #2193b0 0%, #6dd5ed 100%);
        }

        /* Animations */
        @keyframes fadeIn {
            from { opacity: 0; }
            to { opacity: 1; }
        }

        @keyframes slideUp {
            from {
                opacity: 0;
                transform: translateY(20px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }

        @keyframes slideIn {
            from {
                opacity: 0;
                transform: translateX(-20px);
            }
            to {
                opacity: 1;
                transform: translateX(0);
            }
        }

        /* Mobile responsiveness */
        @media (max-width: 768px) {
            [data-testid="stAppViewContainer"] {
                padding: 1rem;
            }

            .title-container {
                padding: 1.5rem 1rem;
            }

            .title-container h1 {
                font-size: 1.8rem;
            }

            .stCard {
                padding: 1rem;
            }

            .metric-value {
                font-size: 1.5rem;
            }
        }
    </style>
""", unsafe_allow_html=True)
def generate_styled_pdf(title: str, content: str, timestamp: str) -> BytesIO:
    """Generate a beautifully styled PDF with the response content."""
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )

    # Styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.HexColor('#1E88E5')
    )

    timestamp_style = ParagraphStyle(
        'Timestamp',
        parent=styles['Normal'],
        fontSize=10,
        textColor=colors.gray
    )

    content_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=12,
        leading=14
    )

    # Build the PDF content
    elements = []

    # Title
    elements.append(Paragraph(title, title_style))
    elements.append(Paragraph(f"Generated on: {timestamp}", timestamp_style))
    elements.append(Spacer(1, 20))

    # Content - split into paragraphs
    paragraphs = content.split('\n\n')
    for para in paragraphs:
        if para.strip():
            elements.append(Paragraph(para, content_style))
            elements.append(Spacer(1, 12))

    # Build the PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

class DocumentProcessor:
    def __init__(self, api_key: str):
        self.client = Groq(api_key=api_key)
        self.logger = logging.getLogger(__name__)

    def process_file(self, file) -> dict:
        """Process a single file and return metadata"""
        start_time = time.time()
        text = self.extract_text(file)
        processing_time = time.time() - start_time

        return {
            'filename': file.name,
            'size': len(file.getvalue()),
            'text_length': len(text),
            'processing_time': processing_time,
            'content': text
        }

    def extract_text(self, file) -> str:
        try:
            if file.type == "application/pdf":
                with pdfplumber.open(file) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(file)
                return "\n".join(paragraph.text for paragraph in doc.paragraphs)
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                ppt = Presentation(file)
                return "\n".join(shape.text for slide in ppt.slides
                               for shape in slide.shapes if hasattr(shape, "text"))
            elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                df = pd.read_excel(file)
                return df.to_string(index=False)
            elif file.type in ["text/plain", "application/octet-stream"]:
                return file.getvalue().decode("utf-8")
            else:
                return "Unsupported file format"
        except Exception as e:
            self.logger.error(f"Error processing file: {str(e)}")
            return f"Error processing file: {str(e)}"

    def process_document(self, content: str, task: str = "summarize",
                        question: Optional[str] = None) -> str:
        try:
            if not content.strip():
                return "No content to process"

            prompt_content = {
                "summarize": f"Summarize the following content concisely:\n\n{content}",
                "ask_question": f"Answer based on the content provided:\n\nContent:\n{content}\n\nQuestion: {question}",
                "combine": f"Combine and sort the following content without losing details:\n\n{content}"
            }.get(task, "Invalid task")

            chat_completion = self.client.chat.completions.create(
                messages=[{"role": "user", "content": prompt_content}],
                model="llama-3.1-70b-versatile",
            )
            return chat_completion.choices[0].message.content

        except Exception as e:
            self.logger.error(f"Error in API call: {str(e)}")
            return f"Error processing request: {str(e)}"

def main():
    # Custom title with gradient background
    st.markdown("""
        <div class="title-container">
            <h1>📄 Smart Document Processor</h1>
            <p>Upload your documents and let AI do the magic</p>
        </div>
    """, unsafe_allow_html=True)

    # Initialize processor
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        st.error("🔑 GROQ_API_KEY environment variable not set")
        return

    processor = DocumentProcessor(api_key)

    # Use columns for layout, but adjust for mobile
    use_columns = True if st.session_state.get('browser_width', 1000) > 768 else False

    if use_columns:
        col1, col2 = st.columns([2, 1])
        upload_container = col1
        stats_container = col2
    else:
        upload_container = st
        stats_container = st

    with upload_container:
        st.markdown("""
            <div class="stCard">
                <h3>📎 Upload Your Multiple Documents Below 👇 </h3>
            </div>
        """, unsafe_allow_html=True)

        uploaded_files = st.file_uploader(
            "",
            type=["pdf", "pptx", "docx", "xlsx", "txt", "py", "js", "html", "java", "cpp"],
            accept_multiple_files=True
        )

    with stats_container:
        if uploaded_files:
            total_files = len(uploaded_files)
            total_size = sum(len(file.getvalue()) for file in uploaded_files)

            metrics_container = st.container()
            with metrics_container:
                if use_columns:
                    mc1, mc2 = st.columns(2)
                    with mc1:
                        st.markdown(f"""
                            <div class="metric-container">
                                <div class="metric-value">{total_files}</div>
                                <div class="metric-label">Files Uploaded</div>
                            </div>
                        """, unsafe_allow_html=True)
                    with mc2:
                        st.markdown(f"""
                            <div class="metric-container">
                                <div class="metric-value">{total_size/1024:.1f}KB</div>
                                <div class="metric-label">Total Size</div>
                            </div>
                        """, unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                        <div class="metric-container">
                            <div class="metric-value">{total_files}</div>
                            <div class="metric-label">Files Uploaded</div>
                        </div>
                        <div class="metric-container">
                            <div class="metric-value">{total_size/1024:.1f}KB</div>
                            <div class="metric-label">Total Size</div>
                        </div>
                    """, unsafe_allow_html=True)

    if uploaded_files:
        st.markdown("<div class='stCard'>", unsafe_allow_html=True)

        # Process files with progress bar
        progress_text = "Processing your documents..."
        progress_bar = st.progress(0)

        processed_files = []
        for idx, file in enumerate(uploaded_files):
            processed = processor.process_file(file)
            processed_files.append(processed)
            progress_bar.progress((idx + 1) / len(uploaded_files))

        combined_text = "\n\n".join(file['content'] for file in processed_files)

        # Task selection
        task = st.radio("🎯 Select Task", ["Summarize", "Ask Questions", "Combine"])

        response = None  # Store the response for PDF generation

        if task == "Ask Questions":
            question_container = st.container()
            with question_container:
                hint_questions = {
                    "HEll":"Recommend the key topics or subjects I should prioritize to effectively prepare for the exam?",
                    "main_points": "What are the main points?",
                    "findings": "What are the key findings?",
                    #"recommendations": "What are the recommendations?",
                    "custom": "Ask your own question..."
                }

                selected_hint = st.selectbox(
                    "💭 Question Type",
                    options=list(hint_questions.keys()),
                    format_func=lambda x: hint_questions[x]
                )

                if selected_hint == "custom":
                    question = st.text_input("🤔 Your Question:")
                else:
                    question = hint_questions[selected_hint]

                if st.button("🚀 Get Answer"):
                    with st.spinner("🔍 Analyzing..."):
                        response = processor.process_document(
                            combined_text,
                            task="ask_question",
                            question=question
                        )
                        st.markdown(f"""
                            <div class="response-container">
                                <h4>Answer:</h4>
                                <p>{response}</p>
                            </div>
                        """, unsafe_allow_html=True)

        elif task == "Summarize":
            if st.button("📝 Generate Summary"):
                with st.spinner("✍️ Creating summary..."):
                    response = processor.process_document(combined_text, task="summarize")
                    st.markdown(f"""
                        <div class="response-container">
                            <h4>Summary:</h4>
                            <p>{response}</p>
                        </div>
                    """, unsafe_allow_html=True)

        elif task == "Combine":
            response = combined_text
            st.markdown(f"""
                <div class="response-container">
                    <h4>Combined Content:</h4>
                    <p>{response}</p>
                </div>
            """, unsafe_allow_html=True)

        # PDF Download button if there's a response
        if response:
            pdf_buffer = generate_styled_pdf(
                f"Document {task} Results",
                response,
                time.strftime("%Y-%m-%d %H:%M:%S")
            )

            st.download_button(
                label="📥 Download as PDF",
                data=pdf_buffer,
                file_name=f"document_{task.lower()}_results.pdf",
                mime="application/pdf",
                key="download_button",
                help="Click to download the results as a PDF file"
            )

        st.markdown("</div>", unsafe_allow_html=True)

if __name__ == "__main__":
    main()
